from dataclasses import dataclass
from decimal import Decimal
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ums_smart_revenue.finance.account_allocation_read import (
    AllocationProvenance,
    account_allocation_disclosure_token,
)
from ums_smart_revenue.finance.bank_reconciliation import (
    MonthBankReconciliationSummary,
)
from ums_smart_revenue.finance.decimal_formatting import decimal_to_api as _decimal_to_api
from ums_smart_revenue.finance.net_revenue import MonthNetRevenueSummary
from ums_smart_revenue.finance.payment_matching import MonthlyPaymentMatchSummary
from ums_smart_revenue.finance.smart_alerts import MonthlySmartAlertSummary
from ums_smart_revenue.reports.exports import ExportJobEntry

BRANDED_SLIDE_NAMES = (
    "Cover",
    "Monthly highlights",
    "Holding revenue summary",
    "Sector comparison",
    "Company comparison",
    "Channel ranking",
    "Revenue deduction explanation",
    "Payment gap analysis",
    "Outside-CMS issues",
    "Action items",
)

_SLIDE_SOURCES = {
    "Cover": "export_job_metadata",
    "Monthly highlights": "computed_finance_summary",
    "Holding revenue summary": "net_revenue_summary",
    "Sector comparison": "channel_registry_and_revenue_facts",
    "Company comparison": "channel_registry_and_revenue_facts",
    "Channel ranking": "monthly_revenue_facts",
    "Revenue deduction explanation": (
        "source_net_revenue_manual_overrides_deduction_components_and_account_allocations"
    ),
    "Payment gap analysis": "adsense_payments_and_bank_reconciliation",
    "Outside-CMS issues": "channel_registry_and_smart_alerts",
    "Action items": "smart_alerts_and_reconciliation_status",
}

_BRAND_BLUE = RGBColor(37, 99, 235)
_BRAND_TEXT = RGBColor(17, 24, 39)
_BRAND_MUTED = RGBColor(75, 85, 99)
_BRAND_LINE = RGBColor(217, 222, 230)


class BrandedSlidePackValidationError(ValueError):
    """Exception raised when validation of a branded slide pack fails."""


@dataclass(frozen=True)
class BrandedSlide:
    """Data class representing a single slide in the branded slide pack."""

    name: str
    source: str
    status: str
    sensitive: bool

    def to_api(self) -> dict[str, object]:
        """Convert the instance to a dictionary formatted for API consumption.

        Returns:
            Dictionary with the instance's name, source, status, and sensitive flag.
        """
        return {
            "name": self.name,
            "source": self.source,
            "status": self.status,
            "sensitive": self.sensitive,
        }


@dataclass(frozen=True)
class BrandedSlidePackReport:
    """Data class representing the full branded slide pack report with
    export job and slide summaries.
    """

    export_job: ExportJobEntry
    slides: tuple[BrandedSlide, ...]
    net_revenue: MonthNetRevenueSummary
    payment_match: MonthlyPaymentMatchSummary
    bank_reconciliation: MonthBankReconciliationSummary
    smart_alerts: MonthlySmartAlertSummary
    account_allocation_provenance: AllocationProvenance

    def to_api(self) -> dict[str, object]:
        """Convert the BrandedSlidePackReport instance into a dictionary
        formatted for API consumption.
        """
        return {
            "export_id": self.export_job.id,
            "export_type": self.export_job.export_type,
            "artifact_type": "BRANDED_FINANCE_SLIDE_PACK",
            "status": "READY_FOR_GENERATION",
            "month": self.export_job.month,
            "currency": self.export_job.currency,
            "scope_type": self.export_job.scope_type,
            "scope_id": self.export_job.scope_id,
            "month_lock_status": self.export_job.month_lock_status,
            "slides": [slide.to_api() for slide in self.slides],
            "executive_summary": _executive_summary(
                export_job=self.export_job,
                net_revenue=self.net_revenue,
                payment_match=self.payment_match,
                bank_reconciliation=self.bank_reconciliation,
                smart_alerts=self.smart_alerts,
            ),
        }


def build_branded_slide_pack_report(
    *,
    export_job: ExportJobEntry,
    net_revenue: MonthNetRevenueSummary,
    payment_match: MonthlyPaymentMatchSummary,
    bank_reconciliation: MonthBankReconciliationSummary,
    smart_alerts: MonthlySmartAlertSummary,
    account_allocation_provenance: AllocationProvenance = AllocationProvenance(
        source="live_compute"
    ),
) -> BrandedSlidePackReport:
    """
    Construct a BrandedSlidePackReport, validating that summary data matches the
    export job and uses USD currency.
    """
    if export_job.export_type != "BRANDED_SLIDE_PACK":
        raise BrandedSlidePackValidationError(
            "branded slide pack only supports BRANDED_SLIDE_PACK exports"
        )
    _validate_same_month(
        export_job=export_job,
        net_revenue=net_revenue,
        payment_match=payment_match,
        bank_reconciliation=bank_reconciliation,
        smart_alerts=smart_alerts,
    )
    if export_job.currency != "USD":
        raise BrandedSlidePackValidationError(
            "branded slide pack requires USD currency"
        )
    if payment_match.currency != "USD" or bank_reconciliation.currency != "USD":
        raise BrandedSlidePackValidationError(
            "branded slide pack requires USD source summaries"
        )

    return BrandedSlidePackReport(
        export_job=export_job,
        slides=tuple(
            BrandedSlide(
                name=name,
                source=_SLIDE_SOURCES[name],
                status="READY",
                sensitive=True,
            )
            for name in BRANDED_SLIDE_NAMES
        ),
        net_revenue=net_revenue,
        payment_match=payment_match,
        bank_reconciliation=bank_reconciliation,
        smart_alerts=smart_alerts,
        account_allocation_provenance=account_allocation_provenance,
    )


def build_branded_slide_pack_pptx(report: BrandedSlidePackReport) -> bytes:
    """
    Generate a PowerPoint presentation as a byte string based on the provided
    BrandedSlidePackReport.
    """
    presentation = Presentation()
    presentation.core_properties.title = "UMS Branded Finance Report"
    presentation.core_properties.subject = f"Finance export {report.export_job.month}"
    presentation.core_properties.author = "UMS Smart Revenue Control Center"
    net_revenue = report.net_revenue
    payment_match = report.payment_match
    bank_reconciliation = report.bank_reconciliation

    _add_cover_slide(presentation, report)
    _add_content_slide(
        presentation,
        "Monthly highlights",
        [
            "Total Net Revenue USD: "
            f"{_decimal_to_api(net_revenue.total_net_revenue_usd)}",
            "Adjusted Gross Revenue USD: "
            f"{_decimal_to_api(net_revenue.total_adjusted_gross_revenue_usd)}",
            f"Payment status: {payment_match.status}",
            f"Bank status: {bank_reconciliation.status}",
        ],
    )
    _add_content_slide(
        presentation,
        "Holding revenue summary",
        [
            f"Calculated channels: {report.net_revenue.calculated_channel_count}",
            f"Total channels: {report.net_revenue.channel_count}",
            f"Month lock status: {report.export_job.month_lock_status}",
        ],
    )
    _add_content_slide(
        presentation,
        "Sector comparison",
        _scope_bullets(report, "sector"),
    )
    _add_content_slide(
        presentation,
        "Company comparison",
        _scope_bullets(report, "company"),
    )
    _add_content_slide(
        presentation,
        "Channel ranking",
        [
            f"{channel.youtube_channel_id}: "
            f"{_decimal_to_api(channel.net_revenue_usd)} USD"
            for channel in sorted(
                report.net_revenue.channels,
                key=lambda item: (
                    item.net_revenue_usd
                    if item.net_revenue_usd is not None
                    else Decimal("-Infinity")
                ),
                reverse=True,
            )[:8]
        ]
        or ["No channel revenue rows are available for this export scope."],
    )
    # ========================================================================
    # Purpose: Render the deduction-explanation slide, surfacing all three month
    #   deduction figures — total, channel-direct, and account-allocated.
    # Database/ORM: None (reads the already-computed MonthNetRevenueSummary).
    # Standards: The channel-direct + account-allocated split SUPPLEMENTS (never
    #   replaces) total_deduction_amount_usd; all three values must stay on the
    #   slide. None values serialize via _decimal_to_api (blank, never "0").
    # Blast Radius: PPTX presentation only — no allocation math, source-of-truth,
    #   auth, audit, or Neo4j impact. Must stay numerically consistent with the
    #   XLSX workbook / PDF deduction-breakdown surfaces and the net-revenue API.
    # Connections:
    #   - File: backend/ums_smart_revenue/finance/net_revenue.py -> Source of the
    #     total_*_deduction_amount_usd aggregates rendered here.
    # ========================================================================
    # The final bullet discloses the account-allocation source (committed
    # snapshot vs live compute / fallback) from the read-switch provenance stored
    # on the report. Read-only presentation; token text owned by
    # account_allocation_disclosure_token. Mirrors the XLSX / PDF disclosure
    # surfaces — no allocation recompute on this path.
    _add_content_slide(
        presentation,
        "Revenue deduction explanation",
        [
            "Total deduction amount USD: "
            f"{_decimal_to_api(net_revenue.total_deduction_amount_usd)}",
            "Channel-direct deduction USD: "
            f"{_decimal_to_api(net_revenue.total_channel_direct_deduction_amount_usd)}",
            "Account-allocated deduction USD: "
            f"{_decimal_to_api(net_revenue.total_account_allocated_deduction_amount_usd)}",
            "Deductions use SQL revenue facts plus approved manual overrides; "
            "the channel-direct portion is sourced from channel-level deduction "
            "components and the account-allocated portion from account-level "
            "deduction allocations.",
            "Pending overrides are shown as risk and do not change net revenue.",
            account_allocation_disclosure_token(
                report.account_allocation_provenance
            ),
        ],
    )
    _add_content_slide(
        presentation,
        "Payment gap analysis",
        [
            "YouTube revenue total USD: "
            f"{_decimal_to_api(payment_match.youtube_revenue_total_usd)}",
            "AdSense paid amount: "
            f"{_decimal_to_api(payment_match.adsense_paid_amount)}",
            f"Payment gap USD: {_decimal_to_api(payment_match.payment_gap_usd)}",
            f"Bank gap USD: {_decimal_to_api(bank_reconciliation.bank_gap_usd)}",
        ],
    )
    _add_content_slide(
        presentation,
        "Outside-CMS issues",
        _outside_cms_bullets(report),
    )
    _add_content_slide(
        presentation,
        "Action items",
        _action_item_bullets(report),
    )

    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _add_cover_slide(
    presentation: Presentation, report: BrandedSlidePackReport
) -> None:
    """Add a cover slide with the branded bar, title, and footer to the presentation."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_brand_bar(slide)
    _add_textbox(
        slide,
        left=Inches(0.7),
        top=Inches(1.5),
        width=Inches(8.8),
        height=Inches(0.7),
        text="UMS Branded Finance Report",
        font_size=Pt(28),
        color=_BRAND_TEXT,
    )
    _add_textbox(
        slide,
        left=Inches(0.72),
        top=Inches(2.45),
        width=Inches(8.6),
        height=Inches(0.5),
        text=(
            f"{report.export_job.month} | {report.export_job.scope_type}"
            + (
                f" | {report.export_job.scope_id}"
                if report.export_job.scope_id is not None
                else ""
            )
        ),
        font_size=Pt(14),
        color=_BRAND_MUTED,
    )
    _add_footer(slide)


def _add_content_slide(
    presentation: Presentation, title: str, bullets: list[str]
) -> None:
    """Add a content slide with a title and bullet points to the presentation."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_brand_bar(slide)
    _add_textbox(
        slide,
        left=Inches(0.55),
        top=Inches(0.55),
        width=Inches(9.0),
        height=Inches(0.45),
        text=title,
        font_size=Pt(20),
        color=_BRAND_TEXT,
    )
    body = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.3), Inches(8.8), Inches(4.9)
    ).text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = _BRAND_TEXT
        paragraph.space_after = Pt(8)
    _add_footer(slide)


def _add_brand_bar(slide) -> None:
    """Add a blue brand bar at the top of the slide for UMS branding."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(0.14),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _BRAND_BLUE
    shape.line.color.rgb = _BRAND_BLUE


def _add_footer(slide) -> None:
    """Add a footer line and text to the bottom of the slide."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(6.82),
        Inches(8.95),
        Inches(0.01),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _BRAND_LINE
    line.line.color.rgb = _BRAND_LINE
    _add_textbox(
        slide,
        left=Inches(0.55),
        top=Inches(6.9),
        width=Inches(5.4),
        height=Inches(0.25),
        text="UMS Smart Revenue Control Center",
        font_size=Pt(8),
        color=_BRAND_MUTED,
    )


def _add_textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    font_size,
    color: RGBColor,
) -> None:
    """Add a textbox with the given text, font size, and color."""
    frame = slide.shapes.add_textbox(left, top, width, height).text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = font_size
    paragraph.font.color.rgb = color


def _scope_bullets(report: BrandedSlidePackReport, scope_name: str) -> list[str]:
    """Generate bullet points summarizing the report scope and net revenue details."""
    return [
        f"{scope_name.title()} scope: {report.export_job.scope_id or 'global'}",
        "Total net revenue USD: "
        f"{_decimal_to_api(report.net_revenue.total_net_revenue_usd)}",
        f"Channels included: {report.net_revenue.channel_count}",
    ]


def _outside_cms_bullets(report: BrandedSlidePackReport) -> list[str]:
    """Generate bullet points for channels missing official net revenue sources."""
    missing_source_count = report.net_revenue.missing_net_source_count
    if missing_source_count == 0:
        return ["No missing source revenue issues were detected in this export scope."]
    return [f"Channels missing official net revenue source: {missing_source_count}"]


def _action_item_bullets(report: BrandedSlidePackReport) -> list[str]:
    """Generate finance action items from report matching statuses."""
    if report.payment_match.status != "PAYMENT_MATCHED":
        return ["Resolve AdSense payment matching before distributing this deck."]
    if report.bank_reconciliation.status != "BANK_CONFIRMED":
        return ["Confirm bank receipt rows before distributing this deck."]
    if report.smart_alerts.alerts:
        return ["Review smart alerts and record follow-up owners."]
    return ["No blocking finance action items detected for this export scope."]


def _validate_same_month(
    *,
    export_job: ExportJobEntry,
    net_revenue: MonthNetRevenueSummary,
    payment_match: MonthlyPaymentMatchSummary,
    bank_reconciliation: MonthBankReconciliationSummary,
    smart_alerts: MonthlySmartAlertSummary,
) -> None:
    """Ensure all report summaries use the same export month; raise error if not."""
    months = {
        export_job.month,
        net_revenue.month,
        payment_match.month,
        bank_reconciliation.month,
        smart_alerts.month,
    }
    if len(months) != 1:
        raise BrandedSlidePackValidationError(
            "branded slide pack source summaries must use the export month"
        )


def _executive_summary(
    *,
    export_job: ExportJobEntry,
    net_revenue: MonthNetRevenueSummary,
    payment_match: MonthlyPaymentMatchSummary,
    bank_reconciliation: MonthBankReconciliationSummary,
    smart_alerts: MonthlySmartAlertSummary,
) -> dict[str, object]:
    """Build an executive summary dict with key report metrics and statuses."""
    return {
        "month": export_job.month,
        "scope_type": export_job.scope_type,
        "scope_id": export_job.scope_id,
        "currency": export_job.currency,
        "month_lock_status": export_job.month_lock_status,
        "net_revenue_status": net_revenue.status,
        "payment_match_status": payment_match.status,
        "bank_reconciliation_status": bank_reconciliation.status,
        "smart_alert_status": smart_alerts.status,
        "smart_alert_count": len(smart_alerts.alerts),
        "total_adjusted_gross_revenue_usd": _decimal_to_api(
            net_revenue.total_adjusted_gross_revenue_usd
        ),
        "total_net_revenue_usd": _decimal_to_api(net_revenue.total_net_revenue_usd),
        "total_deduction_amount_usd": _decimal_to_api(
            net_revenue.total_deduction_amount_usd
        ),
        # Deduction-split totals SUPPLEMENT (never replace) total_deduction_amount_usd
        # and are read-only projections of MonthNetRevenueSummary; see the
        # deduction-slide contract block in build_branded_slide_pack_pptx.
        "total_channel_direct_deduction_amount_usd": _decimal_to_api(
            net_revenue.total_channel_direct_deduction_amount_usd
        ),
        "total_account_allocated_deduction_amount_usd": _decimal_to_api(
            net_revenue.total_account_allocated_deduction_amount_usd
        ),
        "payment_gap_usd": _decimal_to_api(payment_match.payment_gap_usd),
        "bank_gap_usd": _decimal_to_api(bank_reconciliation.bank_gap_usd),
        "channel_count": net_revenue.channel_count,
        "calculated_channel_count": net_revenue.calculated_channel_count,
    }
