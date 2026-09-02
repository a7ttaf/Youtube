"""Tests for export preview API artifacts, downloads, and audit behavior."""

import hashlib
from datetime import UTC, date, datetime
from decimal import Decimal
from io import BytesIO
from types import SimpleNamespace
from typing import NoReturn, cast
from uuid import UUID, uuid4

import pytest
from fastapi import HTTPException
from fastapi.testclient import TestClient
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import create_engine, select
from sqlalchemy.exc import SQLAlchemyError
from sqlalchemy.orm import Session
from starlette.responses import Response as StarletteResponse

from ums_smart_revenue.api.dependencies_audit import current_audit_sink
from ums_smart_revenue.api.exports import (
    _artifact_download_headers,
    _artifact_metadata_audit_details,
    _persist_generated_export_artifact,
    _prepared_artifact_response,
    _require_persisted_artifact_bytes,
)
from ums_smart_revenue.app import create_app
from ums_smart_revenue.auth.audit_service import InMemoryAuditSink
from ums_smart_revenue.db.finance_models import (
    AdSensePaymentORM,
    BankReconciliationEntryORM,
    FinanceBase,
    FinanceMonthCloseORM,
    MonthlyChannelRevenueFactORM,
    RevenueManualOverrideORM,
)
from ums_smart_revenue.db.org_models import (
    ChannelGroupMemberORM,
    ChannelGroupORM,
    OrgBase,
    OrgUnitORM,
    YouTubeChannelORM,
)
from ums_smart_revenue.db.report_models import ExportJobORM, ReportBase
from ums_smart_revenue.db.security_models import AuditLogORM, SecurityBase, UserORM
from ums_smart_revenue.reports.artifact_storage import FileSystemExportArtifactStore
from ums_smart_revenue.reports.exports import (
    ExportJobEntry,
    ExportJobTerminalStateError,
    ExportJobValidationError,
)

SECTOR_ID = UUID("00000000-0000-0000-0000-000000023001")
COMPANY_ID = UUID("00000000-0000-0000-0000-000000023101")
CHANNEL_ROW_ID = UUID("00000000-0000-0000-0000-000000023201")
USER_ID = UUID("00000000-0000-0000-0000-000000023301")
EXPORT_ID = UUID("00000000-0000-0000-0000-000000023401")
GROUP_ID = UUID("00000000-0000-0000-0000-000000023501")
OTHER_USER_ID = UUID("00000000-0000-0000-0000-000000023601")


class _FailingAuditCommitUnitOfWork:
    """SQL-like audit unit of work that fails at the durability boundary."""

    def __init__(self) -> None:
        self.commit_attempted = False
        self.rollback_attempted = False
        self.expunge_attempted = False

    def commit(self) -> None:
        """Raise as a SQL audit commit would after audit rows were appended."""
        self.commit_attempted = True
        raise SQLAlchemyError("audit commit failed")

    def rollback(self) -> None:
        """Record rollback after a failed audit commit."""
        self.rollback_attempted = True

    def expunge_all(self) -> None:
        """Record identity-map cleanup after rollback."""
        self.expunge_attempted = True


class _FailingAuditCommitSink(InMemoryAuditSink):
    """Capture requested audit rows, then fail when the route commits them."""

    def __init__(self) -> None:
        super().__init__()
        self.unit_of_work = _FailingAuditCommitUnitOfWork()
        self.sql_unit_of_work = self.unit_of_work


def _utc_timestamp(value: datetime) -> float:
    """Convert a datetime value to a UTC timestamp (seconds since epoch)."""
    if value.tzinfo is None:
        return value.replace(tzinfo=UTC).timestamp()
    return value.astimezone(UTC).timestamp()


def auth_headers(role: str, scope_type: str = "global") -> dict[str, str]:
    """Generate authentication headers for a user with a given role and scope."""
    return {
        "x-user-id": str(USER_ID),
        "x-user-email": f"{role}@example.com",
        "x-role": role,
        "x-scope-type": scope_type,
        "x-ums-trusted-gateway-token": "pytest-trusted-gateway-token",
    }


def build_database_url(tmp_path) -> str:
    """Build a SQLAlchemy URL for a temporary SQLite database."""
    return f"sqlite+pysqlite:///{(tmp_path / f'{uuid4()}.db').as_posix()}"


def seed_database(
    database_url: str,
    *,
    export_type: str = "FINANCE_EXCEL",
    scope_type: str = "global",
    scope_id: str | None = None,
    requested_by: UUID = USER_ID,
    include_group: bool = False,
    include_previous_fact: bool = False,
) -> None:
    """Seed export preview data, including optional group or prior facts."""
    resolved_scope_id = scope_id
    if resolved_scope_id is None and scope_type == "company":
        resolved_scope_id = str(COMPANY_ID)
    if resolved_scope_id is None and scope_type == "group":
        resolved_scope_id = str(GROUP_ID)

    engine = create_engine(database_url)
    SecurityBase.metadata.create_all(engine)
    OrgBase.metadata.create_all(engine)
    FinanceBase.metadata.create_all(engine)
    ReportBase.metadata.create_all(engine)
    with Session(engine) as session:
        rows = [
            UserORM(
                id=USER_ID,
                email="export-preview@example.com",
                display_name="Export Preview User",
            ),
            OrgUnitORM(
                id=SECTOR_ID,
                parent_id=None,
                type="SECTOR",
                name="TV",
                active=True,
            ),
            OrgUnitORM(
                id=COMPANY_ID,
                parent_id=SECTOR_ID,
                type="COMPANY",
                name="Company A",
                active=True,
            ),
            YouTubeChannelORM(
                id=CHANNEL_ROW_ID,
                youtube_channel_id="channel-tv-a",
                channel_name="Channel TV A",
                primary_org_unit_id=COMPANY_ID,
                cms_status="INSIDE_CMS",
                revenue_required=True,
                active=True,
            ),
            FinanceMonthCloseORM(
                month="2026-03",
                status="LOCKED",
                allocation_rule_payload={},
            ),
            MonthlyChannelRevenueFactORM(
                id=uuid4(),
                month="2026-03",
                youtube_channel_id="channel-tv-a",
                source_kind="YOUTUBE_CMS",
                source_report_id="cms-report-2026-03",
                gross_revenue_usd=Decimal("1000.00"),
                net_revenue_usd=Decimal("880.00"),
                views=250000,
                watch_time_minutes=Decimal("7200.50"),
                confidence_score=Decimal("0.9825"),
                imported_by=USER_ID,
            ),
            RevenueManualOverrideORM(
                id=uuid4(),
                month="2026-03",
                youtube_channel_id="channel-tv-a",
                adjustment_revenue_usd=Decimal("50.00"),
                reason="Approved finance correction",
                status="APPROVED",
                created_by=USER_ID,
                approved_by=uuid4(),
                approved_at=datetime(2026, 4, 24, tzinfo=UTC),
                approval_reason="Approved correction",
            ),
            AdSensePaymentORM(
                id=uuid4(),
                month="2026-03",
                payment_name="March AdSense",
                payment_date=date(2026, 4, 21),
                payment_amount=Decimal("1000.00"),
                payment_currency="USD",
                payment_status="PAID",
                raw_payload={"source": "pytest"},
                source_account_id="pub-1",
                imported_by=USER_ID,
            ),
            BankReconciliationEntryORM(
                id=uuid4(),
                month="2026-03",
                bank_reference="WIRE-2026-03",
                bank_received_date=date(2026, 4, 22),
                bank_received_amount=Decimal("1000.00"),
                bank_received_currency="USD",
                bank_received_amount_usd=Decimal("1000.00"),
                transfer_fee_usd=Decimal("0.00"),
                fx_difference_usd=Decimal("0.00"),
                recorded_by=USER_ID,
            ),
            ExportJobORM(
                id=EXPORT_ID,
                export_type=export_type,
                scope_type=scope_type,
                scope_id=resolved_scope_id,
                month="2026-03",
                currency="USD",
                requested_by=requested_by,
                status="QUEUED",
                month_lock_status="LOCKED",
                include_confidence_notes=True,
                include_manual_override_notes=True,
            ),
        ]
        if include_previous_fact:
            rows.append(
                MonthlyChannelRevenueFactORM(
                    id=uuid4(),
                    month="2026-02",
                    youtube_channel_id="channel-tv-a",
                    source_kind="YOUTUBE_CMS",
                    source_report_id="cms-report-2026-02",
                    gross_revenue_usd=Decimal("100.00"),
                    net_revenue_usd=Decimal("88.00"),
                    views=200000,
                    watch_time_minutes=Decimal("6200.50"),
                    confidence_score=Decimal("0.9800"),
                    imported_by=USER_ID,
                )
            )
        if include_group:
            rows.extend(
                [
                    ChannelGroupORM(
                        id=GROUP_ID,
                        name="March Finance Review",
                        group_type="FINANCE_GROUP",
                        active=True,
                    ),
                    ChannelGroupMemberORM(
                        group_id=GROUP_ID,
                        channel_id=CHANNEL_ROW_ID,
                    ),
                ]
            )
        session.add_all(rows)
        session.commit()


def test_finance_admin_previews_finance_workbook_with_sensitive_audit(tmp_path):
    """Test that finance admins can preview the finance workbook
    and that the audit logs include sensitive events.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM).order_by(AuditLogORM.event_type)).all()

    assert response.status_code == 200
    payload = response.json()
    assert payload["artifact_type"] == "FINANCE_EXCEL_WORKBOOK_PREVIEW"
    assert payload["executive_summary"]["total_net_revenue_usd"] == "930"
    assert payload["executive_summary"]["payment_match_status"] == "PAYMENT_MATCHED"
    assert payload["executive_summary"]["bank_reconciliation_status"] == ("BANK_CONFIRMED")
    assert payload["executive_summary"]["total_channel_direct_deduction_amount_usd"] == "0"
    assert payload["executive_summary"]["total_account_allocated_deduction_amount_usd"] == "0"
    net_revenue_summary = payload["source_summaries"]["net_revenue"]
    assert net_revenue_summary["total_channel_direct_deduction_amount_usd"] == "0"
    assert net_revenue_summary["total_account_allocated_deduction_amount_usd"] == "0"
    assert [sheet["name"] for sheet in payload["sheets"]] == [
        "Executive Summary",
        "Monthly Close",
        "Company Breakdown",
        "Sector Breakdown",
        "Channel Breakdown",
        "Deductions",
        "Payment Gap",
        "Confidence Notes",
        "Raw Appendix",
    ]
    assert {event.event_type for event in audit_events} == {
        "AUDIT_LOG_VIEWED",
        "BANK_RECONCILIATION_VIEWED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }
    assert all(event.sensitive for event in audit_events)


def test_scoped_finance_workbook_omits_month_wide_cash_without_attribution(tmp_path):
    """Test that scoped finance workbook previews omit month-wide
    cash without attribution for company scope.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url, scope_type="company")
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM).order_by(AuditLogORM.event_type)).all()

    assert response.status_code == 200
    payload = response.json()
    assert payload["executive_summary"]["total_net_revenue_usd"] == "930"
    assert payload["executive_summary"]["payment_match_status"] == ("MISSING_ADSENSE_PAYMENT")
    assert payload["executive_summary"]["bank_reconciliation_status"] == ("MISSING_ADSENSE_PAYMENT")
    # Scoped finance-artifact exports now emit PAYMENT_VIEWED alongside
    # REVENUE_VIEWED (export net consumes account allocations); month-wide cash
    # exposure (BANK_RECONCILIATION_VIEWED) stays global-only. Scoped exports
    # do NOT emit AUDIT_LOG_VIEWED because the audit-derived skipped-row signal
    # is suppressed for non-global scopes until source rows can be tied to the
    # export's frozen channel set.
    assert [event.event_type for event in audit_events] == [
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    ]


def test_finance_export_preview_includes_revenue_trend_alerts(tmp_path):
    """Test that finance export preview includes revenue trend
    anomaly alerts when previous month data is present.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url, include_previous_fact=True)
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("finance_admin"),
    )

    assert response.status_code == 200
    alerts = response.json()["source_summaries"]["smart_alerts"]["alerts"]
    trend_alert = next(
        (alert for alert in alerts if alert["code"] == "REVENUE_TREND_ANOMALY"),
        None,
    )
    assert trend_alert is not None
    assert trend_alert["details"]["channels"][0]["youtube_channel_id"] == "channel-tv-a"


def test_artifact_metadata_audit_details_marks_incomplete_metadata():
    """Test that artifact metadata audit details correctly
    identify and mark incomplete metadata fields.
    """
    details = _artifact_metadata_audit_details(
        SimpleNamespace(
            file_url="file-store://exports/export-id/finance.xlsx",
            artifact_filename=None,
            artifact_content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            artifact_byte_size=None,
            artifact_checksum_sha256=None,
        )
    )

    assert details["artifact_locator_redacted"] is True
    assert "file_url" not in details
    assert details["artifact_metadata_complete"] is False
    assert set(details["artifact_metadata_missing_fields"]) == {
        "artifact_filename",
        "artifact_byte_size",
        "artifact_checksum_sha256",
    }
    assert "artifact_filename" not in details


def test_group_scoped_finance_workbook_records_channel_revenue_audit(tmp_path):
    """Test that group-scoped finance workbook preview records
    channel revenue audit events for group scope.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url, scope_type="group", include_group=True)
    client = TestClient(
        create_app(database_url=database_url),
        raise_server_exceptions=False,
    )

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_scopes = session.scalars(
            select(AuditLogORM).order_by(AuditLogORM.event_type, AuditLogORM.id)
        ).all()

    assert response.status_code == 200
    assert ("REVENUE_VIEWED", "channel", "channel-tv-a") in {
        (event.event_type, event.scope_type, event.scope_id) for event in audit_scopes
    }


def test_get_export_returns_not_found_for_missing_group_scope(tmp_path):
    """Legacy group exports without a snapshot return 404 with the group
    lookup error when the source group has been deleted. The caller in this
    test owns the export so the requested_by ownership gate passes and the
    request reaches the permission check.
    """
    database_url = build_database_url(tmp_path)
    seed_database(
        database_url,
        scope_type="group",
        requested_by=USER_ID,
        include_group=False,
    )
    client = TestClient(
        create_app(database_url=database_url),
        raise_server_exceptions=False,
    )

    response = client.get(
        f"/exports/{EXPORT_ID}",
        headers=auth_headers("finance_admin"),
    )

    assert response.status_code == 404
    assert response.json()["detail"] == f"Group not found: {GROUP_ID}"


def test_get_export_returns_not_found_for_other_users_export(tmp_path):
    """Codex P1: GET /exports/{id} must not leak another user's export
    metadata even when the caller has matching scope permissions.
    """
    database_url = build_database_url(tmp_path)
    seed_database(
        database_url,
        scope_type="company",
        requested_by=OTHER_USER_ID,
    )
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}",
        headers=auth_headers("finance_admin"),
    )

    assert response.status_code == 404
    assert response.json()["detail"] == "Export job not found"


def test_export_operator_cannot_preview_finance_workbook(tmp_path):
    """Verify that export operators without finance permission cannot preview finance workbooks."""
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("export_operator"),
    )

    assert response.status_code == 403
    assert response.json()["detail"] == "Missing permission: exports.revenue"


def test_export_operator_cannot_probe_artifact_type_through_finance_preview(
    tmp_path,
):
    """Ensure export operators cannot preview finance workbooks for non-finance export types."""
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="ANALYTICS_SUMMARY_CSV")
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("export_operator"),
    )

    assert response.status_code == 403
    assert response.json()["detail"] == "Missing permission: exports.revenue"


def test_finance_workbook_preview_rejects_analytics_export_job(tmp_path):
    """Ensure finance workbook preview returns 422 when export job is not finance Excel."""
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="ANALYTICS_SUMMARY_CSV")
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook-preview",
        headers=auth_headers("finance_admin"),
    )

    assert response.status_code == 422
    assert response.json()["detail"] == (
        "finance workbook preview only supports FINANCE_EXCEL exports"
    )


def test_finance_admin_downloads_generated_finance_workbook_with_audit(tmp_path):
    """Verify that finance admins can download generated finance
    workbooks and that audit events are recorded.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM).order_by(AuditLogORM.event_type)).all()

    assert response.status_code == 200
    assert response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert "ums-finance-2026-03-global.xlsx" in response.headers["content-disposition"]
    workbook = load_workbook(BytesIO(response.content), data_only=True)
    assert workbook.sheetnames == [
        "Executive Summary",
        "Monthly Close",
        "Company Breakdown",
        "Sector Breakdown",
        "Channel Breakdown",
        "Deductions",
        "Payment Gap",
        "Confidence Notes",
        "Raw Appendix",
    ]
    assert workbook["Executive Summary"]["B2"].value == "2026-03"
    assert workbook["Channel Breakdown"]["A2"].value == "channel-tv-a"
    assert {event.event_type for event in audit_events} == {
        "BANK_RECONCILIATION_VIEWED",
        "EXPORT_DOWNLOADED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }
    assert all(event.sensitive for event in audit_events)


def test_finance_workbook_prepare_is_bodyless_and_real_get_reauthorizes(
    tmp_path,
    monkeypatch,
):
    """Prepare without a download audit, then independently authorize the native GET."""
    artifact_dir = tmp_path / "export-artifacts"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))
    route = f"/exports/{EXPORT_ID}/finance-workbook.xlsx"
    engine = create_engine(database_url)
    original_response_call = StarletteResponse.__call__
    status_when_prepare_response_started: list[str] = []

    async def observe_prepare_commit(response, scope, receive, send):
        """Record the durable-metadata state when the prepare 204 starts."""
        if response.status_code == 204 and scope.get("path") == route:
            with Session(engine) as observation_session:
                observed_job = observation_session.get(ExportJobORM, EXPORT_ID)
                assert observed_job is not None
                status_when_prepare_response_started.append(observed_job.status)
        await original_response_call(response, scope, receive, send)

    monkeypatch.setattr(StarletteResponse, "__call__", observe_prepare_commit)

    denied_prepare = client.get(
        f"{route}?prepare=true",
        headers=auth_headers("export_operator"),
    )
    with Session(engine) as session:
        denied_job = session.get(ExportJobORM, EXPORT_ID)

    assert denied_prepare.status_code == 403
    assert denied_prepare.json()["detail"] == "Missing permission: exports.revenue"
    assert denied_job is not None
    assert denied_job.status == "QUEUED"

    prepared = client.get(
        f"{route}?prepare=true",
        headers=auth_headers("finance_admin"),
    )

    with Session(engine) as session:
        prepared_job = session.get(ExportJobORM, EXPORT_ID)
        preparation_audits = session.scalars(select(AuditLogORM)).all()

    assert prepared.status_code == 204
    assert prepared.content == b""
    assert prepared.headers["cache-control"] == "no-store"
    assert status_when_prepare_response_started == ["COMPLETED"]
    assert prepared_job is not None
    assert prepared_job.status == "COMPLETED"
    assert prepared_job.file_url is not None
    assert prepared_job.artifact_checksum_sha256 is not None
    # FIX (review: prepare generates artifact with no audit trail): the
    # prepare leg READS sensitive finance source data, so it now records its
    # read audit trail — but never EXPORT_DOWNLOADED, which only the
    # delivering GET may claim.
    assert {event.event_type for event in preparation_audits} == {
        "BANK_RECONCILIATION_VIEWED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }

    # The bodyless prepare response is not a bearer grant: the ordinary GET
    # enters the complete route authorization again and still denies a role
    # without finance-export visibility.
    denied = client.get(route, headers=auth_headers("export_operator"))
    assert denied.status_code == 403
    assert denied.json()["detail"] == "Missing permission: exports.revenue"

    downloaded = client.get(route, headers=auth_headers("finance_admin"))
    with Session(engine) as session:
        download_audits = session.scalars(
            select(AuditLogORM).order_by(AuditLogORM.event_type)
        ).all()

    assert downloaded.status_code == 200
    assert downloaded.headers["cache-control"] == "no-store"
    assert "ums-finance-2026-03-global.xlsx" in downloaded.headers["content-disposition"]
    assert hashlib.sha256(downloaded.content).hexdigest() == (prepared_job.artifact_checksum_sha256)
    assert {event.event_type for event in download_audits} == {
        "BANK_RECONCILIATION_VIEWED",
        "EXPORT_DOWNLOADED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }


def test_prepared_artifact_response_fails_closed_when_commit_fails():
    """Do not expose the 204 start signal when artifact metadata is not durable."""

    class FailingCommitSession:
        """Session stand-in whose commit always fails, testing fail-close."""

        @staticmethod
        def commit() -> NoReturn:
            """Raise, simulating a durability failure at metadata commit."""
            raise SQLAlchemyError("commit failed")

    with pytest.raises(HTTPException) as raised:
        _prepared_artifact_response(
            session=cast(Session, FailingCommitSession()), filename="ums-finance.xlsx"
        )

    assert raised.value.status_code == 503
    assert raised.value.detail == "Export artifact persistence unavailable"


def test_prepared_artifact_response_rejects_unsafe_filename_before_committing():
    """The 204 must not announce a download whose filename would fail the header.

    Validation precedes the metadata commit, so an unsafe persisted name
    fails the prepare leg closed without touching the session — the download
    leg can never reach its post-audit header failure.
    """
    commits: list[str] = []

    class RecordingSession:
        """Session stand-in recording its commits."""

        @staticmethod
        def commit() -> None:
            """Record that the prepare leg reached its metadata commit."""
            commits.append("commit")

    with pytest.raises(HTTPException) as raised:
        _prepared_artifact_response(
            session=cast(Session, RecordingSession()), filename='ums-finance-2026-03".xlsx'
        )

    assert raised.value.status_code == 500
    assert raised.value.detail == "Export artifact filename is unsafe"
    assert commits == [], "filename validation must run before the metadata commit"


@pytest.mark.parametrize(
    "poisoned",
    [
        'ums-finance-2026-03".xlsx',
        "ums-finance-2026-03\\.xlsx",
        "ums-finance-2026-03\x00.xlsx",
        "ums-finance-2026-03/../../../etc.xlsx",
        ".",
        "..",
        "  ",
    ],
)
def test_artifact_download_headers_fail_closed_on_unsafe_persisted_filename(poisoned):
    """A persisted filename that could inject or split Content-Disposition fails closed.

    Today's generators write fixed machine names, but the header builder is the
    last choke point every artifact passes through; a future writer must not be
    able to smuggle a quote, backslash, forward slash (a path-shaped name is a
    traversal attempt), control character, or dot-only name into the header.
    """
    with pytest.raises(HTTPException) as raised:
        _artifact_download_headers(poisoned)

    assert raised.value.status_code == 500
    assert raised.value.detail == "Export artifact filename is unsafe"


def test_openapi_documents_the_prepare_parameter_on_every_artifact_route():
    """All four artifact GETs must publish the ``prepare`` handshake parameter."""
    app = create_app(database_url="sqlite+pysqlite:///:memory:")
    with TestClient(app) as client:
        schema = client.get("/openapi.json").json()

    artifact_paths = [
        path
        for path in schema["paths"]
        if path.startswith("/exports/{export_id}/")
        and (
            path.endswith((".csv", ".xlsx", ".pdf", ".pptx"))
        )
    ]
    assert len(artifact_paths) == 4, artifact_paths
    for path in artifact_paths:
        parameters = schema["paths"][path]["get"].get("parameters", [])
        assert any(p.get("name") == "prepare" for p in parameters), (
            f"{path} must document the prepare parameter"
        )


@pytest.mark.parametrize(
    ("export_type", "route", "artifact_prefix"),
    [
        ("FINANCE_EXCEL", f"/exports/{EXPORT_ID}/finance-workbook.xlsx", b"PK"),
        ("EXECUTIVE_PDF", f"/exports/{EXPORT_ID}/executive.pdf", b"%PDF-"),
        ("BRANDED_SLIDE_PACK", f"/exports/{EXPORT_ID}/branded-slide-pack.pptx", b"PK"),
    ],
)
def test_finance_artifact_download_fails_closed_when_audit_commit_fails(
    tmp_path,
    monkeypatch,
    export_type,
    route,
    artifact_prefix,
):
    """Return 503 before finance artifact bytes start when audit commit fails."""
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type=export_type)
    app = create_app(database_url=database_url)
    audit_sink = _FailingAuditCommitSink()
    app.dependency_overrides[current_audit_sink] = lambda: audit_sink
    engine = create_engine(database_url)
    original_response_call = StarletteResponse.__call__
    response_starts: list[int] = []
    response_bodies: list[bytes] = []

    async def observe_response_start(response, scope, receive, send):
        """Wrap ASGI send to capture response starts and body chunks."""
        async def observe_send(message):
            """Forward one ASGI message, recording starts and body chunks."""
            if message["type"] == "http.response.start":
                response_starts.append(message["status"])
            if message["type"] == "http.response.body":
                response_bodies.append(message.get("body", b""))
            await send(message)

        await original_response_call(response, scope, receive, observe_send)

    monkeypatch.setattr(StarletteResponse, "__call__", observe_response_start)

    response = TestClient(app).get(route, headers=auth_headers("finance_admin"))

    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 503
    assert response.json()["detail"] == "Export artifact audit unavailable"
    assert response_starts == [503]
    assert not any(body.startswith(artifact_prefix) for body in response_bodies)
    assert audit_sink.unit_of_work.commit_attempted is True
    assert audit_sink.unit_of_work.rollback_attempted is True
    assert audit_sink.unit_of_work.expunge_attempted is True
    assert {record.event_type for record in audit_sink.records} >= {"EXPORT_DOWNLOADED"}
    assert export_job is not None
    # The tenant session commits BEFORE the audit boundary, so a failed audit
    # commit leaves the job durably COMPLETED over its persisted artifact —
    # a retry serves it and records a real download, and no false
    # EXPORT_DOWNLOADED exists in the durable audit trail.
    assert export_job.status == "COMPLETED"
    assert audit_events == []


class _ValueErrorAuditSink(InMemoryAuditSink):
    """Sink whose transaction context raises a plain ValueError, not a SQL error.

    record_audit_event's own normalization can raise ValueError (a required
    reason, or a downgraded sensitive permission_override); the route's audit
    boundary must contain those exactly like SQLAlchemyError — rollback plus a
    controlled 503 — instead of leaking a raw 500 with a dirty sink session.
    """

    def __init__(self) -> None:
        super().__init__()
        self.unit_of_work = _FailingAuditCommitUnitOfWork()
        self.sql_unit_of_work = self.unit_of_work

    def transaction(self):
        """Raise a non-SQLAlchemy exception before the audit block is entered."""
        raise ValueError("audit reason normalization failed")


@pytest.mark.parametrize(
    ("export_type", "route"),
    [
        ("FINANCE_EXCEL", f"/exports/{EXPORT_ID}/finance-workbook.xlsx"),
        ("EXECUTIVE_PDF", f"/exports/{EXPORT_ID}/executive.pdf"),
        ("BRANDED_SLIDE_PACK", f"/exports/{EXPORT_ID}/branded-slide-pack.pptx"),
    ],
)
def test_finance_artifact_download_fails_closed_when_audit_record_raises_non_sql(
    tmp_path,
    monkeypatch,
    export_type,
    route,
):
    """A non-SQLAlchemy audit exception gets the same rollback + 503 treatment."""
    artifact_dir = tmp_path / "export-artifacts"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type=export_type)
    app = create_app(database_url=database_url)
    audit_sink = _ValueErrorAuditSink()
    app.dependency_overrides[current_audit_sink] = lambda: audit_sink

    response = TestClient(app).get(route, headers=auth_headers("finance_admin"))

    assert response.status_code == 503
    assert response.json()["detail"] == "Export artifact audit unavailable"
    # The non-SQL failure still rolled the sink unit of work back explicitly.
    assert audit_sink.unit_of_work.rollback_attempted is True
    assert audit_sink.unit_of_work.expunge_attempted is True


def test_finance_artifact_download_rejects_unsafe_filename_before_audit(tmp_path, monkeypatch):
    """A poisoned persisted filename fails closed BEFORE any download audit row.

    EXPORT_DOWNLOADED means bytes were delivered. Validating the filename only
    at the header builder — after the audit commit — recorded that event for a
    download that then failed with a raw 500. Validation now runs first.
    """
    artifact_dir = tmp_path / "export-artifacts"
    poisoned_name = "ums-finance-2026-03-global/../../../etc.xlsx"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))
    route = f"/exports/{EXPORT_ID}/finance-workbook.xlsx"

    prepared = client.get(f"{route}?prepare=true", headers=auth_headers("finance_admin"))
    assert prepared.status_code == 204

    # Poison the persisted filename behind storage's back, simulating a row
    # written before the storage boundary's separator/traversal rule existed.
    engine = create_engine(database_url)
    with Session(engine) as session:
        job = session.get(ExportJobORM, EXPORT_ID)
        assert job is not None
        job.artifact_filename = poisoned_name
        session.commit()

    response = client.get(route, headers=auth_headers("finance_admin"))

    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 500
    assert response.json()["detail"] == "Export artifact filename is unsafe"
    # Only the prepare leg's READ audits exist; the GET audited nothing.
    assert "EXPORT_DOWNLOADED" not in {event.event_type for event in audit_events}


def test_finance_workbook_download_persists_artifact_and_completes_job(
    tmp_path,
    monkeypatch,
):
    """Persist the workbook artifact and complete the export job."""
    artifact_dir = tmp_path / "export-artifacts"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))

    # The tenant session's completion metadata must be DURABLE before the
    # artifact bytes start: a failed yield-dependency teardown commit after
    # the response would leave served bytes over unpersisted job state.
    status_when_response_started: list[str | None] = []
    engine_hook = create_engine(database_url)
    original_response_call = StarletteResponse.__call__

    async def observe_response_start(response, scope, receive, send):
        """Record the persisted job status at the moment bytes first start."""

        async def observe_send(message):
            """Capture the job status once, at http.response.start."""
            if message["type"] == "http.response.start" and not status_when_response_started:
                with Session(engine_hook) as snapshot_session:
                    snapshot_job = snapshot_session.get(ExportJobORM, EXPORT_ID)
                status_when_response_started.append(
                    snapshot_job.status if snapshot_job is not None else None
                )
            await send(message)

        await original_response_call(response, scope, receive, observe_send)

    monkeypatch.setattr(StarletteResponse, "__call__", observe_response_start)

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)

    expected_filename = "ums-finance-2026-03-global.xlsx"
    expected_uri = f"file-store://exports/{EXPORT_ID}/{expected_filename}"
    persisted_file = artifact_dir / "exports" / str(EXPORT_ID) / expected_filename
    assert response.status_code == 200
    assert status_when_response_started == ["COMPLETED"]
    assert export_job is not None
    assert export_job.status == "COMPLETED"
    assert export_job.completed_at is not None
    assert export_job.file_url == expected_uri
    assert export_job.artifact_filename == expected_filename
    assert export_job.artifact_content_type == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert export_job.artifact_byte_size == len(response.content)
    assert export_job.artifact_checksum_sha256 == hashlib.sha256(response.content).hexdigest()
    assert export_job.failure_reason is None
    assert persisted_file.read_bytes() == response.content

    cached_response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("finance_admin"),
    )

    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert cached_response.status_code == 200
    assert cached_response.content == response.content
    event_types = [event.event_type for event in audit_events]
    assert event_types.count("EXPORT_DOWNLOADED") == 2
    assert "AUDIT_LOG_VIEWED" not in event_types


def test_finance_workbook_storage_failure_leaves_export_retryable_without_download_audit(
    tmp_path,
    monkeypatch,
):
    """Leave the export retryable when workbook artifact storage fails."""
    # Transient artifact-store outages (here: the configured artifact
    # directory is actually a regular file) must leave the export in its
    # pre-failure status so the caller can retry once storage recovers. The
    # endpoint surfaces 503 but does not terminalize the job — moving it to
    # FAILED would force operators to un-fail before the next attempt.
    artifact_root_file = tmp_path / "not-a-directory"
    artifact_root_file.write_text("not a directory", encoding="utf-8")
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_root_file))
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(
        create_app(database_url=database_url),
        raise_server_exceptions=False,
    )

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 503
    assert response.json()["detail"] == "Export artifact storage unavailable"
    assert export_job is not None
    assert export_job.status == "QUEUED"
    assert export_job.completed_at is None
    assert export_job.file_url is None
    assert export_job.artifact_filename is None
    assert export_job.artifact_content_type is None
    assert export_job.artifact_byte_size is None
    assert export_job.artifact_checksum_sha256 is None
    assert export_job.failure_reason is None
    assert audit_events == []


def test_finance_workbook_download_returns_503_when_persisted_artifact_missing(
    tmp_path,
    monkeypatch,
):
    """A COMPLETED export with an unreadable persisted artifact must return
    503 instead of regenerating bytes that would no longer match the
    persisted checksum recorded in audit metadata. The persisted DB row is
    preserved so an operator can rehydrate the artifact out of band.
    """
    artifact_root_file = tmp_path / "not-a-directory"
    artifact_root_file.write_text("not a directory", encoding="utf-8")
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_root_file))
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    completed_at = datetime(2026, 4, 30, 10, 0, tzinfo=UTC)
    existing_filename = "ums-finance-2026-03-global.xlsx"
    existing_uri = f"file-store://exports/{EXPORT_ID}/{existing_filename}"
    engine = create_engine(database_url)
    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)
        assert export_job is not None
        export_job.status = "COMPLETED"
        export_job.completed_at = completed_at
        export_job.file_url = existing_uri
        export_job.artifact_filename = existing_filename
        export_job.artifact_content_type = (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        export_job.artifact_byte_size = 42
        export_job.artifact_checksum_sha256 = "a" * 64
        session.commit()
    client = TestClient(
        create_app(database_url=database_url),
        raise_server_exceptions=False,
    )

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("finance_admin"),
    )

    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 503
    assert response.json()["detail"] == "Export artifact storage unavailable"
    assert export_job is not None
    assert export_job.status == "COMPLETED"
    assert export_job.completed_at is not None
    assert _utc_timestamp(export_job.completed_at) == _utc_timestamp(completed_at)
    assert export_job.file_url == existing_uri
    assert export_job.artifact_filename == existing_filename
    assert export_job.artifact_content_type == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert export_job.artifact_byte_size == 42
    assert export_job.artifact_checksum_sha256 == "a" * 64
    assert export_job.failure_reason is None
    assert "EXPORT_DOWNLOADED" not in {event.event_type for event in audit_events}


def test_persist_generated_export_artifact_cleans_up_when_completion_fails(tmp_path):
    """Verify that artifacts are cleaned up when artifact completion fails."""

    class FailingCompleteRepository:
        """Repository that raises when artifact completion fails."""

        @staticmethod
        def complete_artifact(**_kwargs: object) -> NoReturn:
            """Raise validation error during artifact completion."""
            raise ExportJobValidationError("completion failed")

    artifact_dir = tmp_path / "artifacts"
    export_job = ExportJobEntry(
        id=str(EXPORT_ID),
        export_type="FINANCE_EXCEL",
        scope_type="global",
        scope_id=None,
        month="2026-03",
        currency="USD",
        requested_by=str(USER_ID),
        status="QUEUED",
        file_url=None,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 30, 10, 0, tzinfo=UTC),
        completed_at=None,
    )

    with pytest.raises(ExportJobValidationError, match="completion failed"):
        _persist_generated_export_artifact(
            repository=FailingCompleteRepository(),
            artifact_store=FileSystemExportArtifactStore(artifact_dir),
            export_job=export_job,
            content=b"generated workbook",
            filename="finance.xlsx",
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    assert not (artifact_dir / "exports" / str(EXPORT_ID) / "finance.xlsx").exists()


def test_persist_generated_export_artifact_rejects_stale_failed_terminal_row(tmp_path):
    """Verify that stale failed terminal export jobs are rejected when persisting artifacts."""
    failed_job = ExportJobEntry(
        id=str(EXPORT_ID),
        export_type="FINANCE_EXCEL",
        scope_type="global",
        scope_id=None,
        month="2026-03",
        currency="USD",
        requested_by=str(USER_ID),
        status="FAILED",
        file_url=None,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 30, 10, 0, tzinfo=UTC),
        completed_at=datetime(2026, 4, 30, 10, 1, tzinfo=UTC),
        failure_reason="artifact storage unavailable",
    )

    class FailedTerminalRepository:
        """Repository stub for a failed terminal export job.

        Simulates interactions for an export job that has already failed,
        rejecting completion attempts and allowing retrieval of the failed job entry.
        """

        @staticmethod
        def complete_artifact(**_kwargs: object) -> NoReturn:
            """Raise terminal state errors when a job is finalized."""
            raise ExportJobTerminalStateError(
                f"Export job {EXPORT_ID} is already in terminal status FAILED"
            )

        @staticmethod
        def get_job(export_id: str) -> ExportJobEntry:
            """Fetch the export job entry by its identifier from the repository."""
            assert export_id == str(EXPORT_ID)
            return failed_job

    artifact_dir = tmp_path / "artifacts"
    export_job = ExportJobEntry(
        id=str(EXPORT_ID),
        export_type="FINANCE_EXCEL",
        scope_type="global",
        scope_id=None,
        month="2026-03",
        currency="USD",
        requested_by=str(USER_ID),
        status="QUEUED",
        file_url=None,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 30, 10, 0, tzinfo=UTC),
        completed_at=None,
    )

    persisted_job, response = _persist_generated_export_artifact(
        repository=FailedTerminalRepository(),
        artifact_store=FileSystemExportArtifactStore(artifact_dir),
        export_job=export_job,
        content=b"generated workbook",
        filename="finance.xlsx",
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    assert persisted_job is None
    assert response is not None
    assert response.status_code == 409
    assert not (artifact_dir / "exports" / str(EXPORT_ID) / "finance.xlsx").exists()


@pytest.mark.parametrize(
    ("export_type", "filename", "content_type"),
    [
        ("EXECUTIVE_PDF", "executive.pdf", "application/pdf"),
        (
            "BRANDED_SLIDE_PACK",
            "slides.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    ],
)
def test_persisted_artifact_bytes_win_after_terminal_race(
    tmp_path,
    export_type: str,
    filename: str,
    content_type: str,
):
    """
    Test that when an artifact already exists, the API handles a terminal-side race
    and preserves the first persisted bytes without overwriting.
    """
    artifact_dir = tmp_path / "artifacts"
    persisted_bytes = b"first-writer-persisted-bytes"
    generated_bytes = b"second-writer-generated-bytes"
    artifact_path = artifact_dir / "exports" / str(EXPORT_ID) / filename
    artifact_path.parent.mkdir(parents=True)
    artifact_path.write_bytes(persisted_bytes)
    file_url = f"file-store://exports/{EXPORT_ID}/{filename}"
    completed_job = ExportJobEntry(
        id=str(EXPORT_ID),
        export_type=export_type,
        scope_type="global",
        scope_id=None,
        month="2026-03",
        currency="USD",
        requested_by=str(USER_ID),
        status="COMPLETED",
        file_url=file_url,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 30, 10, 0, tzinfo=UTC),
        completed_at=datetime(2026, 4, 30, 10, 1, tzinfo=UTC),
        artifact_filename=filename,
        artifact_content_type=content_type,
        artifact_byte_size=len(persisted_bytes),
        artifact_checksum_sha256=hashlib.sha256(persisted_bytes).hexdigest(),
    )
    queued_job = ExportJobEntry(
        id=str(EXPORT_ID),
        export_type=export_type,
        scope_type="global",
        scope_id=None,
        month="2026-03",
        currency="USD",
        requested_by=str(USER_ID),
        status="QUEUED",
        file_url=None,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 30, 10, 0, tzinfo=UTC),
        completed_at=None,
    )

    class CompletedRaceRepository:
        """Repository stub that surfaces a terminal race."""

        @staticmethod
        def complete_artifact(**_kwargs: object) -> NoReturn:
            """Raise when another writer has already completed the job."""
            raise ExportJobTerminalStateError(
                f"Export job {EXPORT_ID} is already in terminal status COMPLETED"
            )

        @staticmethod
        def get_job(export_id: str) -> ExportJobEntry:
            """Retrieve the completed export job entry for the given export_id."""
            assert export_id == str(EXPORT_ID)
            return completed_job

    export_job, response = _persist_generated_export_artifact(
        repository=CompletedRaceRepository(),
        artifact_store=FileSystemExportArtifactStore(artifact_dir),
        export_job=queued_job,
        content=generated_bytes,
        filename=filename,
        content_type=content_type,
    )

    assert response is None
    assert export_job == completed_job
    served_bytes, served_filename, served_content_type = _require_persisted_artifact_bytes(
        export_job=export_job,
        expected_export_type=export_type,
        artifact_store=FileSystemExportArtifactStore(artifact_dir),
    )
    assert served_bytes == persisted_bytes
    assert served_bytes != generated_bytes
    assert served_filename == filename
    assert served_content_type == content_type


def test_export_operator_cannot_download_finance_workbook(tmp_path):
    """Test that export operator cannot download a finance workbook
    and receives a 403 error with no audit events.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url)
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/finance-workbook.xlsx",
        headers=auth_headers("export_operator"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 403
    assert response.json()["detail"] == "Missing permission: exports.revenue"
    assert audit_events == []


def test_finance_admin_downloads_generated_executive_pdf_with_audit(
    tmp_path,
    monkeypatch,
):
    """Test that finance admin can download generated executive
    PDF, triggers audit events, and verifies PDF content.
    """
    artifact_dir = tmp_path / "export-artifacts"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="EXECUTIVE_PDF")
    client = TestClient(create_app(database_url=database_url))
    route = f"/exports/{EXPORT_ID}/executive.pdf"

    prepared = client.get(
        f"{route}?prepare=true",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        preparation_audits = session.scalars(select(AuditLogORM)).all()

    assert prepared.status_code == 204
    assert prepared.content == b""
    assert prepared.headers["cache-control"] == "no-store"
    # The prepare leg reads sensitive finance source data, so it records its
    # read audit trail without EXPORT_DOWNLOADED (the delivering GET's event).
    assert {event.event_type for event in preparation_audits} == {
        "BANK_RECONCILIATION_VIEWED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }

    response = client.get(
        route,
        headers=auth_headers("finance_admin"),
    )

    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM).order_by(AuditLogORM.event_type)).all()

    assert response.status_code == 200
    assert response.headers["cache-control"] == "no-store"
    assert response.headers["content-type"] == "application/pdf"
    assert "ums-executive-2026-03-global.pdf" in response.headers["content-disposition"]
    assert response.content.startswith(b"%PDF-")
    text = _extract_pdf_text(response.content)
    assert "UMS Executive Finance Report" in text
    assert "2026-03" in text
    assert "Total Net Revenue USD" in text
    assert "930" in text
    assert "PAYMENT_MATCHED" in text
    assert "BANK_CONFIRMED" in text
    assert {event.event_type for event in audit_events} == {
        "BANK_RECONCILIATION_VIEWED",
        "EXPORT_DOWNLOADED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }
    assert all(event.sensitive for event in audit_events)
    _assert_persisted_export_artifact(
        database_url=database_url,
        artifact_dir=artifact_dir,
        filename="ums-executive-2026-03-global.pdf",
        content_type="application/pdf",
        content=response.content,
    )


def test_export_operator_cannot_download_executive_pdf(tmp_path):
    """Test that export operator cannot download an executive PDF
    and receives a 403 error with no audit events.
    """
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="EXECUTIVE_PDF")
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/executive.pdf",
        headers=auth_headers("export_operator"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 403
    assert response.json()["detail"] == "Missing permission: exports.revenue"
    assert audit_events == []


def test_finance_admin_downloads_generated_branded_slide_pack_with_audit(
    tmp_path,
    monkeypatch,
):
    """Download a branded slide pack and record sensitive audit events."""
    artifact_dir = tmp_path / "export-artifacts"
    monkeypatch.setenv("UMS_EXPORT_ARTIFACT_DIR", str(artifact_dir))
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="BRANDED_SLIDE_PACK")
    client = TestClient(create_app(database_url=database_url))
    route = f"/exports/{EXPORT_ID}/branded-slide-pack.pptx"

    prepared = client.get(
        f"{route}?prepare=true",
        headers=auth_headers("finance_admin"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        preparation_audits = session.scalars(select(AuditLogORM)).all()

    assert prepared.status_code == 204
    assert prepared.content == b""
    assert prepared.headers["cache-control"] == "no-store"
    # The prepare leg reads sensitive finance source data, so it records its
    # read audit trail without EXPORT_DOWNLOADED (the delivering GET's event).
    assert {event.event_type for event in preparation_audits} == {
        "BANK_RECONCILIATION_VIEWED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }

    response = client.get(
        route,
        headers=auth_headers("finance_admin"),
    )

    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM).order_by(AuditLogORM.event_type)).all()

    assert response.status_code == 200
    assert response.headers["cache-control"] == "no-store"
    assert response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert "ums-branded-2026-03-global.pptx" in response.headers["content-disposition"]
    presentation = Presentation(BytesIO(response.content))
    slide_text = "\n".join(_slide_texts(presentation))
    assert len(presentation.slides) == 10
    assert "UMS Branded Finance Report" in slide_text
    assert "Total Net Revenue USD" in slide_text
    assert "930" in slide_text
    assert "PAYMENT_MATCHED" in slide_text
    assert "BANK_CONFIRMED" in slide_text
    assert {event.event_type for event in audit_events} == {
        "BANK_RECONCILIATION_VIEWED",
        "EXPORT_DOWNLOADED",
        "PAYMENT_VIEWED",
        "REVENUE_VIEWED",
    }
    assert all(event.sensitive for event in audit_events)
    _assert_persisted_export_artifact(
        database_url=database_url,
        artifact_dir=artifact_dir,
        filename="ums-branded-2026-03-global.pptx",
        content_type=("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        content=response.content,
    )


def test_export_operator_cannot_download_branded_slide_pack(tmp_path):
    """Deny branded slide-pack downloads for export operators."""
    database_url = build_database_url(tmp_path)
    seed_database(database_url, export_type="BRANDED_SLIDE_PACK")
    client = TestClient(create_app(database_url=database_url))

    response = client.get(
        f"/exports/{EXPORT_ID}/branded-slide-pack.pptx",
        headers=auth_headers("export_operator"),
    )

    engine = create_engine(database_url)
    with Session(engine) as session:
        audit_events = session.scalars(select(AuditLogORM)).all()

    assert response.status_code == 403
    assert response.json()["detail"] == "Missing permission: exports.revenue"
    assert audit_events == []


def _extract_pdf_text(content: bytes) -> str:
    """Extract and return all text from a PDF file given its content as bytes."""
    reader = PdfReader(BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _slide_texts(presentation: Presentation) -> list[str]:
    """Extract and return text from each slide in a PowerPoint presentation."""
    return [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in presentation.slides
    ]


def _assert_persisted_export_artifact(
    *,
    database_url: str,
    artifact_dir,
    filename: str,
    content_type: str,
    content: bytes,
) -> None:
    """Verify persisted export artifact storage and metadata."""
    expected_uri = f"file-store://exports/{EXPORT_ID}/{filename}"
    persisted_file = artifact_dir / "exports" / str(EXPORT_ID) / filename
    engine = create_engine(database_url)
    with Session(engine) as session:
        export_job = session.get(ExportJobORM, EXPORT_ID)

    assert export_job is not None
    assert export_job.status == "COMPLETED"
    assert export_job.completed_at is not None
    assert export_job.file_url == expected_uri
    assert export_job.artifact_filename == filename
    assert export_job.artifact_content_type == content_type
    assert export_job.artifact_byte_size == len(content)
    assert export_job.artifact_checksum_sha256 == hashlib.sha256(content).hexdigest()
    assert export_job.failure_reason is None
    assert persisted_file.read_bytes() == content
