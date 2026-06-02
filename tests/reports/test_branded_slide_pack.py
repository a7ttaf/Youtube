from datetime import UTC, datetime
from decimal import Decimal
from io import BytesIO
from uuid import uuid4

import pytest
from pptx import Presentation

from ums_smart_revenue.finance.bank_reconciliation import (
    MonthBankReconciliationSummary,
)
from ums_smart_revenue.finance.net_revenue import (
    ChannelNetRevenueSummary,
    MonthNetRevenueSummary,
)
from ums_smart_revenue.finance.payment_matching import MonthlyPaymentMatchSummary
from ums_smart_revenue.finance.smart_alerts import MonthlySmartAlertSummary
from ums_smart_revenue.reports.branded_slide_pack import (
    BRANDED_SLIDE_NAMES,
    BrandedSlidePackValidationError,
    build_branded_slide_pack_pptx,
    build_branded_slide_pack_report,
)
from ums_smart_revenue.reports.exports import ExportJobEntry


def test_branded_slide_pack_report_builds_planned_slide_manifest():
    """Slide pack report builds the correct manifest and executive summary."""
    report = build_branded_slide_pack_report(
        export_job=_export_job(export_type="BRANDED_SLIDE_PACK"),
        net_revenue=_net_revenue_summary(),
        payment_match=_payment_match_summary(status="PAYMENT_MATCHED"),
        bank_reconciliation=_bank_summary(status="BANK_CONFIRMED"),
        smart_alerts=_smart_alert_summary(),
    )

    payload = report.to_api()

    assert payload["artifact_type"] == "BRANDED_FINANCE_SLIDE_PACK"
    assert payload["status"] == "READY_FOR_GENERATION"
    assert [slide["name"] for slide in payload["slides"]] == list(BRANDED_SLIDE_NAMES)
    # Provenance: the deduction explanation slide now surfaces an
    # account-allocated split, so its manifest source must disclose the
    # account-allocation inputs.
    slide_sources = {slide["name"]: slide["source"] for slide in payload["slides"]}
    assert slide_sources["Revenue deduction explanation"] == (
        "source_net_revenue_manual_overrides_and_account_allocations"
    )
    assert payload["executive_summary"]["total_net_revenue_usd"] == "930"
    assert payload["executive_summary"]["payment_match_status"] == "PAYMENT_MATCHED"
    assert payload["executive_summary"]["bank_reconciliation_status"] == (
        "BANK_CONFIRMED"
    )


def test_branded_slide_pack_rejects_non_slide_export_type():
    """Slide pack rejects export jobs that are not BRANDED_SLIDE_PACK."""
    with pytest.raises(BrandedSlidePackValidationError) as exc_info:
        build_branded_slide_pack_report(
            export_job=_export_job(export_type="EXECUTIVE_PDF"),
            net_revenue=_net_revenue_summary(),
            payment_match=_payment_match_summary(status="PAYMENT_MATCHED"),
            bank_reconciliation=_bank_summary(status="BANK_CONFIRMED"),
            smart_alerts=_smart_alert_summary(),
        )

    assert str(exc_info.value) == (
        "branded slide pack only supports BRANDED_SLIDE_PACK exports"
    )


def test_branded_slide_pack_pptx_contains_planned_slides_and_summary_values():
    """Generated PPTX contains all planned slides and summary values."""
    report = build_branded_slide_pack_report(
        export_job=_export_job(export_type="BRANDED_SLIDE_PACK"),
        net_revenue=_net_revenue_summary(),
        payment_match=_payment_match_summary(status="PAYMENT_MATCHED"),
        bank_reconciliation=_bank_summary(status="BANK_CONFIRMED"),
        smart_alerts=_smart_alert_summary(),
    )

    pptx_bytes = build_branded_slide_pack_pptx(report)
    presentation = Presentation(BytesIO(pptx_bytes))
    slide_texts = _slide_texts(presentation)
    combined_text = "\n".join(slide_texts)

    assert pptx_bytes.startswith(b"PK")
    assert len(presentation.slides) == len(BRANDED_SLIDE_NAMES)
    assert "UMS Branded Finance Report" in slide_texts[0]
    assert "2026-03" in slide_texts[0]
    assert "Total Net Revenue USD" in combined_text
    assert "930" in combined_text
    assert "PAYMENT_MATCHED" in combined_text
    assert "BANK_CONFIRMED" in combined_text
    for slide_name in BRANDED_SLIDE_NAMES[1:]:
        assert slide_name in combined_text


def test_branded_slide_pack_pptx_handles_missing_channel_net_revenue():
    """PPTX gracefully renders channels with missing net revenue source."""
    report = build_branded_slide_pack_report(
        export_job=_export_job(export_type="BRANDED_SLIDE_PACK"),
        net_revenue=_net_revenue_summary(include_missing_channel=True),
        payment_match=_payment_match_summary(status="PAYMENT_MATCHED"),
        bank_reconciliation=_bank_summary(status="BANK_CONFIRMED"),
        smart_alerts=_smart_alert_summary(),
    )

    pptx_bytes = build_branded_slide_pack_pptx(report)
    presentation = Presentation(BytesIO(pptx_bytes))
    combined_text = "\n".join(_slide_texts(presentation))

    assert pptx_bytes.startswith(b"PK")
    assert "channel-missing-source" in combined_text


def _slide_texts(presentation: Presentation) -> list[str]:
    """Extract all text from each slide in a Presentation."""
    return [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in presentation.slides
    ]


def _export_job(*, export_type: str) -> ExportJobEntry:
    """Build a minimal ExportJobEntry for slide-pack tests."""
    return ExportJobEntry(
        id=str(uuid4()),
        export_type=export_type,
        scope_type="company",
        scope_id="company-a",
        month="2026-03",
        currency="USD",
        requested_by=str(uuid4()),
        status="QUEUED",
        file_url=None,
        month_lock_status="LOCKED",
        include_confidence_notes=True,
        include_manual_override_notes=True,
        created_at=datetime(2026, 4, 1, tzinfo=UTC),
        completed_at=None,
    )


def _net_revenue_summary(
    *, include_missing_channel: bool = False
) -> MonthNetRevenueSummary:
    """Build a MonthNetRevenueSummary for slide-pack tests."""
    channel = ChannelNetRevenueSummary(
        month="2026-03",
        youtube_channel_id="channel-tv-a",
        status="CALCULATED",
        primary_source_kind="YOUTUBE_CMS",
        baseline_gross_revenue_usd=Decimal("1000.00"),
        baseline_net_revenue_usd=Decimal("880.00"),
        approved_manual_override_total_usd=Decimal("50.00"),
        adjusted_gross_revenue_usd=Decimal("1050.00"),
        net_revenue_usd=Decimal("930.00"),
        deduction_amount_usd=Decimal("120.00"),
        channel_direct_deduction_amount_usd=None,
        account_allocated_deduction_amount_usd=None,
        deduction_percentage=Decimal("11.4286"),
        confidence="B_RECONCILED",
        approved_manual_override_count=1,
        pending_manual_override_count=0,
        issues=[],
    )
    channels = [channel]
    if include_missing_channel:
        channels.append(
            ChannelNetRevenueSummary(
                month="2026-03",
                youtube_channel_id="channel-missing-source",
                status="NET_REVENUE_SOURCE_MISSING",
                primary_source_kind=None,
                baseline_gross_revenue_usd=Decimal("100.00"),
                baseline_net_revenue_usd=None,
                approved_manual_override_total_usd=Decimal("0.00"),
                adjusted_gross_revenue_usd=Decimal("100.00"),
                net_revenue_usd=None,
                deduction_amount_usd=None,
                channel_direct_deduction_amount_usd=None,
                account_allocated_deduction_amount_usd=None,
                deduction_percentage=None,
                confidence="E_MISSING",
                approved_manual_override_count=0,
                pending_manual_override_count=0,
                issues=[],
            )
        )
    return MonthNetRevenueSummary(
        month="2026-03",
        status="CALCULATED",
        channel_count=len(channels),
        calculated_channel_count=1,
        missing_net_source_count=1 if include_missing_channel else 0,
        pending_manual_override_count=0,
        total_adjusted_gross_revenue_usd=Decimal("1050.00"),
        total_net_revenue_usd=Decimal("930.00"),
        total_deduction_amount_usd=Decimal("120.00"),
        total_channel_direct_deduction_amount_usd=Decimal("0.00"),
        total_account_allocated_deduction_amount_usd=Decimal("0.00"),
        unallocated_account_deduction_total_usd=None,
        unallocated_account_issues=None,
        channels=channels,
    )


def _payment_match_summary(*, status: str) -> MonthlyPaymentMatchSummary:
    """Build a MonthlyPaymentMatchSummary for slide-pack tests."""
    return MonthlyPaymentMatchSummary(
        month="2026-03",
        currency="USD",
        status=status,
        youtube_revenue_total_usd=Decimal("930.00"),
        adsense_paid_amount=Decimal("930.00"),
        payment_gap_usd=Decimal("0.00"),
        youtube_source_channel_count=1,
        missing_youtube_source_channel_count=0,
        payment_count=1,
        paid_payment_count=1,
        non_paid_payment_count=0,
        unsupported_payment_currency_count=0,
        tolerance_usd=Decimal("0.01"),
        issues=[],
    )


def _bank_summary(*, status: str) -> MonthBankReconciliationSummary:
    """Build a MonthBankReconciliationSummary for slide-pack tests."""
    return MonthBankReconciliationSummary(
        month="2026-03",
        currency="USD",
        status=status,
        adsense_paid_amount_usd=Decimal("930.00"),
        bank_received_amount_usd=Decimal("930.00"),
        bank_gap_usd=Decimal("0.00"),
        transfer_fee_usd=Decimal("0.00"),
        fx_difference_usd=Decimal("0.00"),
        payment_count=1,
        paid_payment_count=1,
        non_paid_payment_count=0,
        unsupported_payment_currency_count=0,
        entry_count=1,
        tolerance_usd=Decimal("0.01"),
        issues=[],
        entries=[],
    )


def _smart_alert_summary() -> MonthlySmartAlertSummary:
    """Build a MonthlySmartAlertSummary for slide-pack tests."""
    return MonthlySmartAlertSummary(
        month="2026-03",
        status="CLEAR",
        highest_severity=None,
        alerts=[],
    )


def _net_revenue_summary_with_breakdown() -> MonthNetRevenueSummary:
    """Build a summary with a COMPONENT_DERIVED channel carrying a real split."""
    channel = ChannelNetRevenueSummary(
        month="2026-03",
        youtube_channel_id="channel-tv-a",
        status="COMPONENT_DERIVED",
        primary_source_kind="ADSENSE",
        baseline_gross_revenue_usd=Decimal("1000.00"),
        baseline_net_revenue_usd=None,
        approved_manual_override_total_usd=Decimal("0.00"),
        adjusted_gross_revenue_usd=Decimal("1000.00"),
        net_revenue_usd=Decimal("870.00"),
        deduction_amount_usd=Decimal("130.00"),
        channel_direct_deduction_amount_usd=Decimal("30.00"),
        account_allocated_deduction_amount_usd=Decimal("100.00"),
        deduction_percentage=Decimal("13.0000"),
        confidence="D_ESTIMATED",
        approved_manual_override_count=0,
        pending_manual_override_count=0,
        issues=[],
    )
    return MonthNetRevenueSummary(
        month="2026-03",
        status="CALCULATED",
        channel_count=1,
        calculated_channel_count=1,
        missing_net_source_count=0,
        pending_manual_override_count=0,
        total_adjusted_gross_revenue_usd=Decimal("1000.00"),
        total_net_revenue_usd=Decimal("870.00"),
        total_deduction_amount_usd=Decimal("130.00"),
        total_channel_direct_deduction_amount_usd=Decimal("30.00"),
        total_account_allocated_deduction_amount_usd=Decimal("100.00"),
        unallocated_account_deduction_total_usd=None,
        unallocated_account_issues=None,
        channels=[channel],
    )


def test_branded_slide_pack_renders_deduction_breakdown_bullets():
    """PPTX deduction slide shows total + channel-direct + account-allocated bullets."""
    report = build_branded_slide_pack_report(
        export_job=_export_job(export_type="BRANDED_SLIDE_PACK"),
        net_revenue=_net_revenue_summary_with_breakdown(),
        payment_match=_payment_match_summary(status="PAYMENT_MATCHED"),
        bank_reconciliation=_bank_summary(status="BANK_CONFIRMED"),
        smart_alerts=_smart_alert_summary(),
    )
    combined_text = "\n".join(
        _slide_texts(Presentation(BytesIO(build_branded_slide_pack_pptx(report))))
    )

    assert "Total deduction amount USD: 130" in combined_text
    assert "Channel-direct deduction USD: 30" in combined_text
    assert "Account-allocated deduction USD: 100" in combined_text
    # Provenance accuracy: the slide must attribute each split portion to its
    # true source — channel-direct to channel-level deduction components,
    # account-allocated to account-level deduction allocations.
    assert "channel-level deduction components" in combined_text
    assert "account-level deduction allocations" in combined_text

    payload = report.to_api()
    assert payload["executive_summary"]["total_channel_direct_deduction_amount_usd"] == "30"
    assert payload["executive_summary"]["total_account_allocated_deduction_amount_usd"] == "100"
